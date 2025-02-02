import os
import re
import json
import sys
import logging
import datetime
from typing import Optional, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import win32com.client  # For converting PPTX to PDF
from natsort import natsorted
from PIL import Image, ImageDraw, ImageFont

# Constants
FONT_PATH = r"C:\Windows\Fonts\times.ttf"  # Path to Times New Roman font file
CM_TO_INCH = 1 / 2.54  # Conversion factor from centimeters to inches

# Photo dimensions and positions (converted from centimeters to inches)
PHOTO_WIDTH = Inches(9.8 * CM_TO_INCH)
PHOTO_HEIGHT = Inches(13.06 * CM_TO_INCH)
PHOTO_TOP = Inches(4.72 * CM_TO_INCH)

# Fixed horizontal positions for 3 columns (in inches)
FIXED_PHOTO_LEFTS = [
    Inches(1.32 * CM_TO_INCH),
    Inches(12.04 * CM_TO_INCH),
    Inches(22.63 * CM_TO_INCH)
]

# Label dimensions and margins (in inches)
LABEL_HEIGHT = Inches(2.0 * CM_TO_INCH)
LABEL_MARGIN_LEFT = Inches(0.25 * CM_TO_INCH)
LABEL_MARGIN_RIGHT = Inches(0.25 * CM_TO_INCH)
LABEL_MARGIN_TOP = Inches(0.13 * CM_TO_INCH)
LABEL_MARGIN_BOTTOM = Inches(0.13 * CM_TO_INCH)
LABEL_FONT_FAMILY = "Times New Roman"
LABEL_FONT_SIZE = 28  # Maximum font size in points

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s"
)
logger = logging.getLogger(__name__)


def extract_number(text: str) -> int:
    """
    Extracts a leading number from a string for sorting.
    """
    match = re.match(r'^(\d+)', text)
    return int(match.group(1)) if match else 999999


def extract_photo_number(filename: str) -> int:
    """
    Extracts a leading number from a filename for sorting.
    """
    match = re.match(r'^(\d+)', filename)
    return int(match.group(1)) if match else 999999


def extract_date_from_folder(folder_name: str) -> Optional[datetime.datetime]:
    """
    Extracts a date in DD.MM.YYYY format from the folder name.
    """
    match = re.search(r"(\d{2}\.\d{2}\.\d{4})", folder_name)
    if match:
        try:
            return datetime.datetime.strptime(match.group(1), "%d.%m.%Y")
        except Exception as e:
            logger.warning("Date parsing error from '%s': %s", folder_name, e)
    return None


def shrink_text_to_fit(text: str, available_width_pt: float, available_height_pt: float,
                       font_path: str, max_font_size: int, min_font_size: int = 8) -> int:
    """
    Iteratively decreases the font size until the text fits within the available width and height.
    Uses Pillow's textbbox method for measuring text dimensions.
    """
    image = Image.new("RGB", (int(available_width_pt), int(available_height_pt)))
    draw = ImageDraw.Draw(image)
    for size in range(max_font_size, min_font_size - 1, -1):
        try:
            font = ImageFont.truetype(font_path, size)
        except IOError:
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        if text_width <= available_width_pt and text_height <= available_height_pt:
            return size
    return min_font_size


# -----------------------------
# Load configuration from JSON file
# -----------------------------
CONFIG_FILE = "config.json"
try:
    with open(CONFIG_FILE, "r", encoding="utf-8") as config_file:
        config = json.load(config_file)
except Exception as e:
    logger.error("Error loading configuration: %s", e)
    sys.exit(1)

folder_path = config.get("folder_path")
output_path = config.get("output_path")
title_content = config.get("title_content")

if not folder_path or not output_path or not title_content:
    logger.error("Missing one or more required configuration parameters.")
    sys.exit(1)


def add_title_slide(prs: Presentation, title_text: str) -> None:
    """
    Adds a title slide to the presentation.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Create a full-slide textbox for the title
    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    # Set title font properties
    p.font.name = 'Times New Roman'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    # Center text vertically
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    logger.info("Title slide added.")


def add_photo_with_label(slide, photo_path: str, left, photo_top, label_text: str) -> None:
    """
    Adds a photo with a header label above it.
    The header text is resized to optimally fit the designated area.
    """
    label_top = photo_top - LABEL_HEIGHT
    textbox = slide.shapes.add_textbox(left, label_top, PHOTO_WIDTH, LABEL_HEIGHT)
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = LABEL_MARGIN_LEFT
    tf.margin_right = LABEL_MARGIN_RIGHT
    tf.margin_top = LABEL_MARGIN_TOP
    tf.margin_bottom = LABEL_MARGIN_BOTTOM

    # Calculate available width and height in points (1 inch = 72 pt)
    available_width_pt = (9.8 * CM_TO_INCH * 72) - ((0.25 * CM_TO_INCH * 72) * 2)
    available_height_pt = (2.0 * CM_TO_INCH * 72) - ((0.13 * CM_TO_INCH * 72) * 2)
    optimal_font_size = shrink_text_to_fit(label_text, available_width_pt, available_height_pt,
                                           FONT_PATH, LABEL_FONT_SIZE)
    logger.info("For text '%s', optimal font size: %d pt", label_text, optimal_font_size)

    p = tf.paragraphs[0]
    p.text = label_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = LABEL_FONT_FAMILY
    p.font.size = Pt(optimal_font_size)

    # Add the photo to the slide
    slide.shapes.add_picture(photo_path, left, photo_top, width=PHOTO_WIDTH, height=PHOTO_HEIGHT)


def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> None:
    """
    Converts a PPTX file to PDF using win32com.
    """
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)  # 32 corresponds to ppSaveAsPDF
        presentation.Close()
        powerpoint.Quit()
        logger.info("PDF saved: %s", pdf_path)
    except Exception as e:
        logger.error("Error converting PPTX to PDF: %s", e)


def get_unique_photos_in_folder(folder: str) -> List[str]:
    """
    Returns a naturally sorted list of unique image files (containing '_stamped')
    from the given folder. For each group (based on the prefix before '_'),
    the longest filename is chosen.
    """
    files = [
        f for f in os.listdir(folder)
        if f.lower().endswith(('.jpg', '.jpeg', '.png')) and "_stamped" in f.lower()
    ]
    groups: dict[str, str] = {}
    for f in files:
        key = f.split("_")[0]
        current = groups.get(key)
        if current is None or len(f) > len(current):
            groups[key] = f
    unique_files = [os.path.join(folder, f) for f in groups.values()]
    return natsorted(unique_files, key=lambda x: extract_photo_number(os.path.basename(x)))


def get_leaf_folders(root: str) -> List[str]:
    """
    Recursively finds and returns leaf folders (directories with no subdirectories).
    """
    leaf_folders = []
    for dirpath, dirnames, _ in os.walk(root):
        if not dirnames:
            leaf_folders.append(dirpath)
    return leaf_folders


def create_photo_report_presentation(root_path: str, output_path: str) -> None:
    """
    Creates a photo report presentation and converts it to PDF.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, title_content)

    leaf_folders = get_leaf_folders(root_path)
    # Sort folders by extracted date; folders without a date are treated as the oldest
    sorted_folders = sorted(
        leaf_folders,
        key=lambda f: extract_date_from_folder(os.path.basename(f)) or datetime.datetime.min
    )
    logger.info("Found leaf folders: %s", sorted_folders)

    all_photos = []
    for folder in sorted_folders:
        # Skip folders that contain unwanted keywords
        if "Лишнее" in folder or "Запас" in folder:
            continue
        work_name = os.path.basename(folder)
        folder_date = extract_date_from_folder(work_name) or datetime.datetime.min
        photos = get_unique_photos_in_folder(folder)
        if not photos:
            logger.warning("No suitable photos found in folder '%s'.", folder)
            continue
        for photo in photos:
            photo_num = extract_photo_number(os.path.basename(photo))
            all_photos.append((folder_date, photo_num, photo, work_name))
    all_photos.sort(key=lambda x: (x[0], x[1]))
    logger.info("Total photos collected: %d", len(all_photos))

    # Calculate the number of slides (3 photos per slide)
    num_slides = (len(all_photos) + 2) // 3
    for slide_idx in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Calculate vertical centering for the photo group (label + photo)
        group_height = LABEL_HEIGHT + PHOTO_HEIGHT
        group_top = (prs.slide_height - group_height) / 2
        photo_top = group_top + LABEL_HEIGHT

        # Get current group of photos (up to 3)
        group = all_photos[slide_idx * 3: slide_idx * 3 + 3]
        n = len(group)
        positions = []
        if n == 3:
            positions = FIXED_PHOTO_LEFTS
        elif n == 2:
            gap = Inches(0.5)
            total_width = 2 * PHOTO_WIDTH + gap
            start_left = (prs.slide_width - total_width) / 2
            positions = [start_left, start_left + PHOTO_WIDTH + gap]
        elif n == 1:
            positions = [(prs.slide_width - PHOTO_WIDTH) / 2]

        for i, (_, _, photo_path, work_name) in enumerate(group):
            add_photo_with_label(slide, photo_path, positions[i], photo_top, work_name)
        logger.info("Created slide %d with %d photos.", slide_idx + 1, n)

    # Ensure the output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    try:
        prs.save(output_path)
        logger.info("Presentation saved: %s", output_path)
    except Exception as e:
        logger.error("Error saving presentation: %s", e)
        sys.exit(1)

    pdf_path = os.path.splitext(output_path)[0] + ".pdf"
    try:
        convert_pptx_to_pdf(output_path, pdf_path)
    except Exception as e:
        logger.error("Error converting to PDF: %s", e)


def main() -> None:
    """
    Main entry point for creating the photo report presentation.
    """
    logger.info("Starting presentation creation...")
    create_photo_report_presentation(folder_path, output_path)
    logger.info("Process completed.")


if __name__ == "__main__":
    main()
