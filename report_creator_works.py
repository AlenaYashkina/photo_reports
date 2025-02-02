import os
import re
import json
import sys
import logging
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import win32com.client  # For converting PPTX to PDF
from PIL import Image, ImageDraw, ImageFont

# -----------------------------
# Constants and Parameters
# -----------------------------
FONT_PATH = r"C:\Windows\Fonts\times.ttf"  # Path to Times New Roman font file

# Configure logger
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s"
)
logger = logging.getLogger(__name__)

CM_TO_INCH = 1 / 2.54

# Photo sizes in inches (converted from centimeters)
PHOTO_WIDTH = Inches(9.8 * CM_TO_INCH)    # ~9.8 cm
PHOTO_HEIGHT = Inches(13.06 * CM_TO_INCH)  # ~13.06 cm

# Label dimensions and margins in inches
LABEL_HEIGHT = Inches(2.0 * CM_TO_INCH)  # ~2 cm
LABEL_MARGIN_LEFT = Inches(0.25 * CM_TO_INCH)
LABEL_MARGIN_RIGHT = Inches(0.25 * CM_TO_INCH)
LABEL_MARGIN_TOP = Inches(0.13 * CM_TO_INCH)
LABEL_MARGIN_BOTTOM = Inches(0.13 * CM_TO_INCH)
LABEL_FONT_FAMILY = "Times New Roman"
LABEL_FONT_SIZE = 28  # Maximum font size for label

# Gap between photos on slide
GAP = Inches(0.5)

# Slide dimensions (e.g., standard 13.33 x 7.5 inches)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


# -----------------------------
# Helper Functions
# -----------------------------
def extract_folder_number(folder_name: str) -> int:
    """
    Extracts the numeric prefix from a folder name for sorting.
    E.g., "01 Example" returns 1.
    """
    match = re.match(r'^(\d+)', folder_name)
    return int(match.group(1)) if match else 999999


def extract_folder_label(folder_name: str) -> str:
    """
    Returns the folder name without its numeric prefix.
    E.g., "01 Example" becomes "Example".
    """
    return re.sub(r'^\d+\s*', '', folder_name)


def shrink_text_to_fit(text: str, available_width_pt: float, available_height_pt: float,
                       font_path: str, max_font_size: int, min_font_size: int = 8) -> int:
    """
    Reduces the font size until the text fits within the available space.
    Uses Pillow to measure text dimensions.
    """
    # Create a dummy image to measure text size
    image = Image.new("RGB", (int(available_width_pt), int(available_height_pt)))
    draw = ImageDraw.Draw(image)
    for size in range(max_font_size, min_font_size - 1, -1):
        try:
            font = ImageFont.truetype(font_path, size)
        except IOError:
            font = ImageFont.load_default()
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        if text_width <= available_width_pt and text_height <= available_height_pt:
            return size
    return min_font_size


def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> None:
    """
    Converts a PPTX file to PDF using win32com.
    """
    powerpoint = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        presentation.SaveAs(pdf_path, 32)  # 32 corresponds to ppSaveAsPDF
        presentation.Close()
        logger.info("PDF saved: %s", pdf_path)
    except Exception as e:
        logger.error("Error converting PPTX to PDF: %s", e)
    finally:
        if powerpoint:
            powerpoint.Quit()


def get_leaf_folders(root: str) -> List[str]:
    """
    Recursively finds and returns leaf folders (folders without subdirectories).
    """
    leaf_folders = []
    for dirpath, dirnames, _ in os.walk(root):
        if not dirnames:
            leaf_folders.append(dirpath)
    return leaf_folders


def get_photos_in_folder(folder: str) -> List[str]:
    """
    Returns a list of image file paths in a folder that contain '_stamped' in their name.
    Supported formats: .jpg, .jpeg, .png.
    """
    files = [
        f for f in os.listdir(folder)
        if f.lower().endswith(('.jpg', '.jpeg', '.png')) and "_stamped" in f.lower()
    ]
    files_sorted = sorted(files)
    return [os.path.join(folder, f) for f in files_sorted]


def emu_to_points(emu: float) -> float:
    """
    Converts a measurement from EMU to points.
    (1 inch = 914400 EMU, 1 inch = 72 pt)
    """
    return emu / 914400 * 72


# -----------------------------
# Presentation Creation Functions
# -----------------------------
def add_title_slide(prs: Presentation, title_text: str) -> None:
    """
    Adds a title slide with centered text.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Times New Roman"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    logger.info("Title slide added.")


def add_slide_with_photos(prs: Presentation, photos: List[str], label_text: str) -> None:
    """
    Adds a slide with a header label and up to three photos.
    The label is placed above the photo area.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Calculate block width for 3 photos and 2 gaps
    block_width = 3 * PHOTO_WIDTH + 2 * GAP
    block_left = (prs.slide_width - block_width) / 2

    # Calculate positions for label and photo group
    group_height = LABEL_HEIGHT + PHOTO_HEIGHT
    group_top = (prs.slide_height - group_height) / 2
    label_top = group_top

    # Add label textbox
    textbox = slide.shapes.add_textbox(block_left, label_top, block_width, LABEL_HEIGHT)
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = LABEL_MARGIN_LEFT
    tf.margin_right = LABEL_MARGIN_RIGHT
    tf.margin_top = LABEL_MARGIN_TOP
    tf.margin_bottom = LABEL_MARGIN_BOTTOM

    # Calculate available width and height in points for the text
    available_width_pt = emu_to_points(block_width) - (
        emu_to_points(LABEL_MARGIN_LEFT) + emu_to_points(LABEL_MARGIN_RIGHT)
    )
    available_height_pt = emu_to_points(LABEL_HEIGHT) - (
        emu_to_points(LABEL_MARGIN_TOP) + emu_to_points(LABEL_MARGIN_BOTTOM)
    )
    optimal_font_size = shrink_text_to_fit(label_text, available_width_pt,
                                           available_height_pt, FONT_PATH,
                                           LABEL_FONT_SIZE)
    logger.info("For label '%s', selected font size: %d pt", label_text, optimal_font_size)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label_text
    run.font.name = LABEL_FONT_FAMILY
    run.font.size = Pt(optimal_font_size)

    # Calculate position for photos (below the label)
    photos_top = group_top + LABEL_HEIGHT
    n = len(photos)
    current_block_width = n * PHOTO_WIDTH + (n - 1) * GAP
    photos_left = (prs.slide_width - current_block_width) / 2

    for i, photo_path in enumerate(photos):
        left = photos_left + i * (PHOTO_WIDTH + GAP)
        slide.shapes.add_picture(photo_path, left, photos_top, width=PHOTO_WIDTH, height=PHOTO_HEIGHT)
    logger.info("Slide created with %d photo(s) and label '%s'.", n, label_text)


def create_photo_report_presentation(root_path: str, output_path: str, title_content: str) -> None:
    """
    Creates a presentation with a title slide and slides for each leaf folder.
    For each leaf folder:
      - Skips folders with names containing "Лишнее" or "Запас".
      - Uses photos with '_stamped' in their names.
      - Groups photos in sets of 3 with the folder label as header.
    After saving the presentation, it converts the file to PDF.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs, title_content)

    leaf_folders = get_leaf_folders(root_path)
    sorted_folders = sorted(leaf_folders, key=lambda f: extract_folder_number(os.path.basename(f)))
    logger.info("Found leaf folders: %s", sorted_folders)

    for folder in sorted_folders:
        folder_name = os.path.basename(folder)
        if "Лишнее" in folder_name or "Запас" in folder_name:
            continue

        photos = get_photos_in_folder(folder)
        if not photos:
            logger.warning("No photos with '_stamped' found in folder '%s'.", folder)
            continue

        folder_label = extract_folder_label(folder_name)

        # Group photos in sets of 3 per slide
        for i in range(0, len(photos), 3):
            group = photos[i:i + 3]
            add_slide_with_photos(prs, group, folder_label)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    try:
        prs.save(output_path)
        logger.info("Presentation saved: %s", output_path)
    except Exception as e:
        logger.error("Error saving presentation: %s", e)
        sys.exit(1)

    pdf_path = os.path.splitext(output_path)[0] + ".pdf"
    try:
        convert_pptx_to_pdf(output_path, pdf_path)
    except Exception as e:
        logger.error("Error converting to PDF: %s", e)


def main() -> None:
    """
    Main function to generate the presentation from a configuration file.
    """
    CONFIG_FILE = "config.json"
    try:
        with open(CONFIG_FILE, "r", encoding="utf-8") as config_file:
            config = json.load(config_file)
    except Exception as e:
        logger.error("Error loading configuration: %s", e)
        sys.exit(1)

    folder_path = config.get("folder_path")
    output_path = config.get("output_path")
    title_content = config.get("title_content")

    if not folder_path or not output_path or not title_content:
        logger.error("One or more required configuration parameters are missing.")
        sys.exit(1)

    logger.info("Starting presentation creation...")
    create_photo_report_presentation(folder_path, output_path, title_content)
    logger.info("Process completed.")


if __name__ == "__main__":
    main()
